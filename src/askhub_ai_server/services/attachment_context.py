from __future__ import annotations

import io
import logging
from dataclasses import dataclass
from uuid import UUID

from sqlalchemy.orm import Session

from askhub_ai_server.core.config import Settings
from askhub_ai_server.models import ChatSession, UserFile
from askhub_ai_server.schemas.chat import ChatAttachment
from askhub_ai_server.services.exceptions import ServiceError
from askhub_ai_server.services.file_storage import FileStorage, get_file_storage

logger = logging.getLogger(__name__)

BEDROCK_DOCUMENT_BLOCK_MAX_BYTES = 3_300_000
PPTX_CONTENT_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

IMAGE_CONTENT_TYPES = {"image/png", "image/jpeg", "image/jpg", "image/gif", "image/webp"}
DOCUMENT_CONTENT_TYPES = {
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
}


@dataclass(frozen=True)
class AttachmentContext:
    text: str
    attachments: list[ChatAttachment]


class AttachmentContextBuilder:
    def __init__(
        self,
        db: Session,
        settings: Settings,
        storage: FileStorage | None = None,
    ) -> None:
        self._db = db
        self._settings = settings
        self._storage = storage

    def build(self, file_ids: list[UUID], session: ChatSession) -> AttachmentContext:
        if not file_ids:
            return AttachmentContext(text="", attachments=[])
        if len(file_ids) > self._settings.max_files_per_message:
            raise ServiceError(400, "too many files attached to a single message")

        parts: list[str] = []
        attachments: list[ChatAttachment] = []
        for file_id in file_ids:
            user_file = self._db.get(UserFile, file_id)
            if user_file is None:
                raise ServiceError(404, f"attached file not found: {file_id}")
            self._authorize_file(user_file, session)
            self._append_file_context(user_file, parts, attachments)

        return AttachmentContext(text="\n\n".join(parts), attachments=attachments)

    def _append_file_context(
        self,
        user_file: UserFile,
        parts: list[str],
        attachments: list[ChatAttachment],
    ) -> None:
        content_type = _normalize_content_type(user_file.content_type)
        try:
            if _is_image_content_type(content_type):
                self._append_image_context(user_file, content_type, parts, attachments)
            elif _is_pptx_content_type(content_type):
                self._append_pptx_context(user_file, parts)
            elif _is_document_content_type(content_type):
                self._append_document_context(user_file, content_type, parts, attachments)
            elif _is_text_context_content_type(content_type, self._settings):
                self._append_text_context(user_file, parts)
            else:
                parts.append(
                    "[첨부 파일 metadata]\n"
                    f"filename={user_file.filename}\n"
                    f"content_type={content_type or 'unknown'}\n"
                    f"file_size={user_file.file_size or 0}"
                )
        except ServiceError:
            raise
        except UnicodeDecodeError as exc:
            raise ServiceError(
                400,
                f"attached file is not valid UTF-8 text: {user_file.id}",
            ) from exc
        except OSError as exc:
            logger.warning("파일 읽기 실패: %s", user_file.storage_path, exc_info=True)
            raise ServiceError(
                500,
                f"attached file content is not available: {user_file.id}",
            ) from exc
        except Exception as exc:
            logger.warning("파일 저장소 읽기 실패: %s", user_file.id, exc_info=True)
            raise ServiceError(
                500,
                f"attached file content is not available: {user_file.id}",
            ) from exc

    def _append_image_context(
        self,
        user_file: UserFile,
        content_type: str,
        parts: list[str],
        attachments: list[ChatAttachment],
    ) -> None:
        content = self._file_storage.read_bytes(user_file, self._settings.max_upload_bytes + 1)
        if len(content) > self._settings.max_upload_bytes:
            raise ServiceError(
                413,
                f"attached file exceeds the configured size limit: {user_file.id}",
            )
        attachments.append(
            ChatAttachment(filename=user_file.filename, content_type=content_type, data=content)
        )
        parts.append(f"[첨부 이미지: {user_file.filename} ({content_type})]")

    def _append_pptx_context(self, user_file: UserFile, parts: list[str]) -> None:
        content = self._file_storage.read_bytes(user_file, self._settings.max_upload_bytes + 1)
        if len(content) > self._settings.max_upload_bytes:
            raise ServiceError(
                413,
                f"attached file exceeds the configured size limit: {user_file.id}",
            )
        parts.append(f"[첨부 파일: {user_file.filename}]\n{_extract_pptx_text(content)}")

    def _append_document_context(
        self,
        user_file: UserFile,
        content_type: str,
        parts: list[str],
        attachments: list[ChatAttachment],
    ) -> None:
        content = self._file_storage.read_bytes(user_file, BEDROCK_DOCUMENT_BLOCK_MAX_BYTES + 1)
        if len(content) > BEDROCK_DOCUMENT_BLOCK_MAX_BYTES:
            raise ServiceError(
                413,
                (
                    "attached document exceeds the Bedrock document block"
                    f" size limit (~3.3 MB): {user_file.id}"
                ),
            )
        attachments.append(
            ChatAttachment(filename=user_file.filename, content_type=content_type, data=content)
        )
        parts.append(f"[첨부 문서: {user_file.filename} ({content_type})]")

    def _append_text_context(self, user_file: UserFile, parts: list[str]) -> None:
        content = self._file_storage.read_bytes(user_file, self._settings.max_file_context_bytes)
        text = content.decode("utf-8")
        parts.append(f"[첨부 파일: {user_file.filename}]\n{text}")

    def _authorize_file(self, user_file: UserFile, session: ChatSession) -> None:
        if user_file.user_id != session.user_id:
            raise ServiceError(404, "attached file not found")
        if session.team_id is not None and user_file.team_id != session.team_id:
            raise ServiceError(404, "attached file not found")
        if user_file.session_id is not None and user_file.session_id != session.id:
            raise ServiceError(403, "file session mismatch")
        if user_file.purpose != "chat_attachment":
            raise ServiceError(400, "only chat attachment files can be attached to chat messages")

    @property
    def _file_storage(self) -> FileStorage:
        if self._storage is None:
            self._storage = get_file_storage(self._settings)
        return self._storage


def _normalize_content_type(content_type: str | None) -> str:
    return content_type.split(";")[0].strip().lower() if content_type else ""


def _is_text_context_content_type(content_type: str, settings: Settings) -> bool:
    return (
        content_type.startswith("text/") or content_type in settings.allowed_upload_content_types
    ) and not _is_image_content_type(content_type)


def _is_image_content_type(content_type: str) -> bool:
    return content_type in IMAGE_CONTENT_TYPES


def _is_document_content_type(content_type: str) -> bool:
    return content_type in DOCUMENT_CONTENT_TYPES


def _is_pptx_content_type(content_type: str) -> bool:
    return content_type == PPTX_CONTENT_TYPE


def _extract_pptx_text(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            parts.append(f"[슬라이드 {i}]\n" + "\n".join(texts))
    return "\n\n".join(parts)
