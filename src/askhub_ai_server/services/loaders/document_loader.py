"""문서 텍스트 추출 — markitdown 기반 마크다운 변환 + 레거시 fallback."""

from __future__ import annotations

import io
import logging
from dataclasses import dataclass, field
from pathlib import Path

logger = logging.getLogger(__name__)

# markitdown으로 변환 가능한 확장자 (바이너리이지만 텍스트 추출 가능)
CONVERTIBLE_EXTENSIONS: set[str] = {
    ".pdf", ".docx", ".pptx", ".xlsx", ".xls", ".csv", ".html",
    ".png", ".jpg", ".jpeg", ".gif", ".webp",  # OCR 대상
}

# 바이너리 파일 제외 확장자 (변환 불가능한 순수 바이너리)
BINARY_EXTENSIONS: set[str] = {
    ".bmp", ".ico", ".svg",
    ".mp3", ".mp4", ".avi", ".mov", ".wav", ".flac",
    ".zip", ".tar", ".gz", ".bz2", ".7z", ".rar",
    ".exe", ".dll", ".so", ".dylib", ".bin",
    ".woff", ".woff2", ".ttf", ".eot", ".otf",
    ".pyc", ".class", ".o", ".obj",
}


@dataclass(frozen=True)
class LoadedPage:
    """추출된 문서 페이지 텍스트."""

    page_number: int
    content: str


@dataclass(frozen=True)
class LoadedDocument:
    """추출된 문서 텍스트."""

    file_path: str
    content: str
    file_type: str
    title: str | None = None
    pages: list[LoadedPage] = field(default_factory=list)


def is_binary_file(file_path: str) -> bool:
    """확장자 기반으로 바이너리 파일 여부를 판별한다."""
    ext = Path(file_path).suffix.lower()
    return ext in BINARY_EXTENSIONS


def extract_pdf_text(data: bytes) -> str:
    """PDF 바이트에서 텍스트를 추출한다 (pymupdf)."""
    return "\n\n".join(page.content for page in extract_pdf_pages(data))


def extract_pdf_pages(data: bytes) -> list[LoadedPage]:
    """PDF 바이트에서 페이지별 텍스트를 추출한다 (pymupdf)."""
    import fitz  # pymupdf

    doc = fitz.open(stream=data, filetype="pdf")
    pages: list[LoadedPage] = []
    for page in doc:
        text = page.get_text().strip()
        if text:
            pages.append(LoadedPage(page_number=page.number + 1, content=text))
    doc.close()
    return pages


def extract_docx_text(data: bytes) -> str:
    """DOCX 바이트에서 텍스트를 추출한다 (python-docx)."""
    from docx import Document

    doc = Document(io.BytesIO(data))
    parts: list[str] = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)
    return "\n\n".join(parts)


def extract_pptx_text(data: bytes) -> str:
    """PPTX 바이트에서 텍스트를 추출한다 (python-pptx)."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            parts.append(f"[슬라이드 {i}]\n" + "\n".join(texts))
    return "\n\n".join(parts)


def extract_text_from_bytes(data: bytes, file_path: str) -> str:
    """파일 확장자에 따라 적절한 텍스트 추출 방법을 선택한다."""
    ext = Path(file_path).suffix.lower()

    if ext == ".pdf":
        return extract_pdf_text(data)
    if ext == ".docx":
        return extract_docx_text(data)
    if ext == ".pptx":
        return extract_pptx_text(data)

    # 일반 텍스트
    return data.decode("utf-8", errors="replace")


def load_document_from_file(file_path: str) -> LoadedDocument | None:
    """로컬 파일 경로에서 문서를 로드하고 텍스트를 추출한다."""
    from askhub_ai_server.services.chunker import detect_file_type

    path = Path(file_path)
    if not path.is_file():
        logger.warning("파일 없음: %s", file_path)
        return None
    if is_binary_file(file_path):
        logger.debug("바이너리 파일 건너뜀: %s", file_path)
        return None

    data = path.read_bytes()
    if not data:
        return None

    pages: list[LoadedPage] = []
    if path.suffix.lower() == ".pdf":
        pages = extract_pdf_pages(data)
        content = "\n\n".join(page.content for page in pages)
    else:
        content = extract_text_from_bytes(data, file_path)
    if not content.strip():
        return None

    return LoadedDocument(
        file_path=file_path,
        content=content,
        file_type=detect_file_type(file_path),
        title=path.name,
        pages=pages,
    )


def load_document_from_bytes(data: bytes, file_path: str) -> LoadedDocument | None:
    """바이트 데이터에서 문서를 로드하고 텍스트를 추출한다."""
    from askhub_ai_server.services.chunker import detect_file_type

    if is_binary_file(file_path):
        return None
    if not data:
        return None

    ext = Path(file_path).suffix.lower()
    pages: list[LoadedPage] = []
    if ext == ".pdf":
        pages = extract_pdf_pages(data)
        content = "\n\n".join(page.content for page in pages)
    else:
        content = extract_text_from_bytes(data, file_path)
    if not content.strip():
        return None

    return LoadedDocument(
        file_path=file_path,
        content=content,
        file_type=detect_file_type(file_path),
        title=Path(file_path).name,
        pages=pages,
    )


# ---------------------------------------------------------------------------
# markitdown 기반 마크다운 변환
# ---------------------------------------------------------------------------


def convert_to_markdown(data: bytes, file_path: str) -> str:
    """markitdown으로 문서/이미지 바이트를 마크다운으로 변환한다.

    지원 포맷: PDF, DOCX, PPTX, XLSX, CSV, HTML, 이미지(OCR).
    변환 실패 시 레거시 텍스트 추출로 fallback한다.
    """
    ext = Path(file_path).suffix.lower()

    try:
        from markitdown import MarkItDown

        md = MarkItDown()
        result = md.convert_stream(io.BytesIO(data), file_extension=ext)
        text = result.text_content
        if text and text.strip():
            return text
    except Exception:
        logger.warning("markitdown 변환 실패, fallback 시도: %s", file_path, exc_info=True)

    # fallback: 기존 텍스트 추출
    if ext == ".pdf":
        return extract_pdf_text(data)
    if ext == ".docx":
        return extract_docx_text(data)
    if ext == ".pptx":
        return extract_pptx_text(data)
    if ext in CONVERTIBLE_EXTENSIONS:
        # XLSX, 이미지 등 fallback 불가 → 빈 문자열
        return ""

    return data.decode("utf-8", errors="replace")


def load_document_as_markdown(data: bytes, file_path: str) -> LoadedDocument | None:
    """문서를 마크다운으로 변환하여 LoadedDocument를 반환한다.

    모든 문서 포맷을 markitdown으로 마크다운 변환한 뒤 content에 저장한다.
    변환 불가능한 바이너리 파일은 None을 반환한다.
    """
    from askhub_ai_server.services.chunker import detect_file_type

    ext = Path(file_path).suffix.lower()

    # 변환 가능한 확장자이면 마크다운 변환 시도
    if ext in CONVERTIBLE_EXTENSIONS:
        if not data:
            return None
        content = convert_to_markdown(data, file_path)
        if not content.strip():
            return None
        return LoadedDocument(
            file_path=file_path,
            content=content,
            file_type=detect_file_type(file_path),
            title=Path(file_path).name,
        )

    # 순수 바이너리 파일은 처리 불가
    if is_binary_file(file_path):
        return None
    if not data:
        return None

    # 텍스트 계열 파일은 그대로 UTF-8 디코딩
    content = data.decode("utf-8", errors="replace")
    if not content.strip():
        return None

    return LoadedDocument(
        file_path=file_path,
        content=content,
        file_type=detect_file_type(file_path),
        title=Path(file_path).name,
    )
